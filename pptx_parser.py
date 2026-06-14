import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt


# ── Shape text extraction ─────────────────────────────────────────────────

def extract_shape_text(shape) -> str:
    """
    Extract all text from a shape.
    Handles: text boxes, titles, grouped shapes (recursive).
    Preserves bullet hierarchy with indentation.
    """
    lines = []

    # Grouped shapes — recurse into children
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child_shape in shape.shapes:
            child_text = extract_shape_text(child_shape)
            if child_text:
                lines.append(child_text)
        return '\n'.join(lines)

    # Text-bearing shapes
    if not shape.has_text_frame:
        return ""

    for para in shape.text_frame.paragraphs:
        para_text = para.text.strip()
        if not para_text:
            continue

        # Indent based on bullet level
        level = para.level or 0
        if level == 0:
            prefix = "• "
        elif level == 1:
            prefix = "  ◦ "
        else:
            prefix = "    ▪ "

        lines.append(f"{prefix}{para_text}")

    return '\n'.join(lines)


# ── Table extraction ──────────────────────────────────────────────────────

def extract_table_from_shape(shape) -> str:
    """
    Extract table content as formatted text.
    First row treated as header.
    Merged cells deduplicated.
    """
    if not shape.has_table:
        return ""

    table = shape.table
    lines = []

    for row_idx, row in enumerate(table.rows):
        cells = []
        prev_text = None

        for cell in row.cells:
            cell_text = cell.text.strip()

            # Deduplicate merged cells (same text repeated)
            if cell_text != prev_text:
                cells.append(cell_text)
            prev_text = cell_text

        if not any(cells):
            continue  # skip empty rows

        row_text = ' | '.join(cells)

        if row_idx == 0:
            lines.append(f"Header: {row_text}")
        else:
            lines.append(f"Row {row_idx}: {row_text}")

    return '\n'.join(lines)


# ── Single slide extraction ───────────────────────────────────────────────

def extract_slide(slide, slide_num: int) -> dict:
    """
    Extract all content from one slide.
    Returns a dict with text content and metadata flags.

    Content extracted:
    - Title (from title placeholder)
    - All text shapes (with bullet hierarchy)
    - All tables (formatted as rows)
    - Speaker notes
    """
    title_text = ""
    body_parts = []
    table_parts = []
    has_tables = False
    has_notes = False

    # NEW — wraps placeholder access in try/except
    for shape in slide.shapes:
        if shape.has_text_frame:
            try:
                ph_fmt = shape.placeholder_format
                if ph_fmt is not None and ph_fmt.idx == 0:
                    title_text = shape.text_frame.text.strip()
                    break
            except ValueError:
                # shape.placeholder_format raises ValueError on non-placeholders
                # Skip and continue looking for title
                continue

    # Fallback: if no title placeholder, use first text shape
    if not title_text:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                title_text = shape.text.strip()[:100]  # truncate long fallbacks
                break

    # ── Extract all shapes ────────────────────────────────────────────────
    # NEW
    for shape in slide.shapes:
        # Skip the title placeholder we already captured
        try:
            ph_fmt = shape.placeholder_format
            if ph_fmt is not None and ph_fmt.idx == 0:
                continue  # this is the title — already captured above
        except ValueError:
            pass  # not a placeholder — process normally
        
        # Tables
        if shape.has_table:
            table_text = extract_table_from_shape(shape)
            if table_text:
                table_parts.append(f"[Table]\n{table_text}")
                has_tables = True
            continue
        
        # Text shapes
        shape_text = extract_shape_text(shape)
        if shape_text:
            body_parts.append(shape_text)

    # ── Extract speaker notes ─────────────────────────────────────────────
    notes_text = ""
    try:
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_text = notes_frame.text.strip()
            if notes_text:
                has_notes = True
    except Exception:
        pass  # some slides have malformed notes — skip gracefully

    # ── Assemble slide chunk ──────────────────────────────────────────────
    chunk_parts = []

    # Header line
    if title_text:
        chunk_parts.append(f"Slide {slide_num}: {title_text}")
    else:
        chunk_parts.append(f"Slide {slide_num}")

    # Body content
    if body_parts:
        chunk_parts.append('\n'.join(body_parts))

    # Tables
    if table_parts:
        chunk_parts.append('\n'.join(table_parts))

    # Speaker notes — clearly labelled
    if notes_text:
        chunk_parts.append(f"[Speaker Notes]\n{notes_text}")

    slide_text = '\n\n'.join(chunk_parts)

    return {
        "text": slide_text,
        "slide_number": slide_num,
        "title": title_text or f"Slide {slide_num}",
        "has_tables": has_tables,
        "has_notes": has_notes
    }


# ── Main PPTX parser ──────────────────────────────────────────────────────

def parse_pptx(file_path: str, filename: str) -> dict:
    """
    Parse a PowerPoint file into slide chunks.

    Returns:
    - slides: list of slide dicts (text + metadata per slide)
    - presentation_meta: title, author, total slides
    - summary_chunk: all slide titles concatenated for overview queries
    """
    try:
        prs = Presentation(file_path)
    except Exception as e:
        return {"error": f"Could not open PPTX file: {e}"}

    total_slides = len(prs.slides)

    if total_slides == 0:
        return {"error": "Presentation has no slides"}

    # ── Presentation metadata ─────────────────────────────────────────────
    props = prs.core_properties
    meta_parts = [f"Presentation: {filename}"]

    if props.title:
        meta_parts.append(f"Title: {props.title}")
    if props.author:
        meta_parts.append(f"Author: {props.author}")
    if props.created:
        meta_parts.append(f"Created: {props.created.strftime('%Y-%m-%d')}")

    meta_parts.append(f"Total slides: {total_slides}")
    presentation_meta_text = '\n'.join(meta_parts)

    # ── Extract all slides ────────────────────────────────────────────────
    slides_data = []

    for i, slide in enumerate(prs.slides, start=1):
        slide_data = extract_slide(slide, i)
        if slide_data["text"].strip():
            slides_data.append(slide_data)

    if not slides_data:
        return {"error": "No content could be extracted from slides"}

    # ── Summary chunk — all titles ────────────────────────────────────────
    # "What topics does this presentation cover?" answers from this chunk
    title_lines = [
        f"Slide {s['slide_number']}: {s['title']}"
        for s in slides_data
    ]
    summary_text = (
        f"Presentation overview — {filename}\n"
        f"Total slides: {total_slides}\n\n"
        "Slide titles:\n" + '\n'.join(title_lines)
    )

    return {
        "filename": filename,
        "type": "pptx",
        "total_slides": total_slides,
        "slides_extracted": len(slides_data),
        "presentation_meta": presentation_meta_text,
        "summary": summary_text,
        "slides": slides_data,
        "slides_with_tables": sum(1 for s in slides_data if s["has_tables"]),
        "slides_with_notes": sum(1 for s in slides_data if s["has_notes"])
    }